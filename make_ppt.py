"""
智水清源路演PPT — 动画+渐变版
python-pptx + 底层XML渐变/动画
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy

# ── 配色 ──
C_BLUE_DARK  = "003399"
C_BLUE       = "0057D9"
C_BLUE_MID   = "0077FF"
C_BLUE_LIGHT = "00AAFF"
C_WHITE      = "FFFFFF"
C_TEXT       = "1A1A2E"
C_TEXT2      = "5A6178"
C_TEXT3      = "9AA0B4"
C_LIGHT_BG   = "EEF6FF"
C_RED        = "FF3B30"
C_ORANGE     = "FF9500"
C_GREEN      = "34C759"

def rgb(s):
    r,g,b = int(s[0:2],16),int(s[2:4],16),int(s[4:6],16)
    return RGBColor(r,g,b)

BLUE_DARK  = rgb(C_BLUE_DARK)
BLUE       = rgb(C_BLUE)
BLUE_MID   = rgb(C_BLUE_MID)
BLUE_LIGHT = rgb(C_BLUE_LIGHT)
WHITE      = rgb(C_WHITE)
TEXT       = rgb(C_TEXT)
TEXT2      = rgb(C_TEXT2)
TEXT3      = rgb(C_TEXT3)
LIGHT_BG   = rgb(C_LIGHT_BG)

# ── 全局设置 ──
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ══════════════════════════════════════════
# XML 渐变辅助
# ══════════════════════════════════════════
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"

def set_shape_gradient(shape, color1, color2, angle=5400000):
    """设置形状线性渐变填充（角度单位：1/60000度）"""
    sp = shape._sp
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(sp, qn("p:spPr"))

    # 清除现有填充
    for old in spPr.findall(qn("a:gradFill")):
        spPr.remove(old)
    # 移除 solidFill
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    # 移除 noFill
    for old in spPr.findall(qn("a:noFill")):
        spPr.remove(old)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"), attrib={"rotWithShape": "1"})
    gs = etree.SubElement(gradFill, qn("a:gsLst"))
    g1 = etree.SubElement(gs, qn("a:gs"), attrib={"pos": "0"})
    etree.SubElement(g1, qn("a:srgbClr"), attrib={"val": color1})
    g2 = etree.SubElement(gs, qn("a:gs"), attrib={"pos": "100000"})
    etree.SubElement(g2, qn("a:srgbClr"), attrib={"val": color2})
    etree.SubElement(gradFill, qn("a:lin"), attrib={
        "ang": str(angle), "scaled": "1"
    })

def set_slide_gradient_bg(slide, color1, color2, color3=None, color4=None):
    """设置幻灯片背景渐变"""
    bg = slide.background
    fill = bg.fill
    fill._xPr.set("gradFill", "")
    grad = etree.SubElement(fill._xPr, qn("a:gradFill"), attrib={"rotWithShape": "1"})
    gs = etree.SubElement(grad, qn("a:gsLst"))

    colors = [color1, color2] + ([color3] if color3 else []) + ([color4] if color4 else [])
    steps = [0, 50000] + ([75000] if color3 else []) + ([100000] if color4 else [])
    for c, pos in zip(colors, steps):
        g = etree.SubElement(gs, qn("a:gs"), attrib={"pos": str(pos)})
        etree.SubElement(g, qn("a:srgbClr"), attrib={"val": c})
    etree.SubElement(grad, qn("a:lin"), attrib={"ang": "5400000", "scaled": "1"})

def set_bg_solid(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill._xPr.set("solidFill", "")
    sf = etree.SubElement(fill._xPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color})

def set_slide_bg_blend(slide, color1, color2, color3="EEF6FF"):
    """三层渐变背景"""
    bg = slide.background
    fill = bg.fill
    # 清除
    for child in list(fill._xPr):
        fill._xPr.remove(child)
    grad = etree.SubElement(fill._xPr, qn("a:gradFill"), attrib={"rotWithShape": "1"})
    gs = etree.SubElement(grad, qn("a:gsLst"))
    for val, pos in [(color1,"0"),(color2,"50000"),(color3,"100000")]:
        g = etree.SubElement(gs, qn("a:gs"), attrib={"pos": pos})
        etree.SubElement(g, qn("a:srgbClr"), attrib={"val": val})
    etree.SubElement(grad, qn("a:lin"), attrib={"ang": "5400000", "scaled": "1"})

def set_rect_gradient(rect, color1, color2, angle=0):
    """给 rect 形状设置渐变"""
    sp = rect._sp
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for old in spPr.findall(qn("a:gradFill")):
        spPr.remove(old)
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    for old in spPr.findall(qn("a:noFill")):
        spPr.remove(old)
    gradFill = etree.SubElement(spPr, qn("a:gradFill"), attrib={"rotWithShape": "1"})
    gs = etree.SubElement(gradFill, qn("a:gsLst"))
    for val, pos in [(color1,"0"),(color2,"100000")]:
        g = etree.SubElement(gs, qn("a:gs"), attrib={"pos": pos})
        etree.SubElement(g, qn("a:srgbClr"), attrib={"val": val})
    etree.SubElement(gradFill, qn("a:lin"), attrib={
        "ang": str(angle), "scaled": "1"
    })

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, round=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_round_rect(slide, l, t, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT,
             italic=False, font_name="PingFang SC"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color) if isinstance(color, str) else color
    run.font.name = font_name
    return txBox

def _c(color):
    """统一转hex字符串"""
    if isinstance(color, str):
        return color
    if isinstance(color, RGBColor):
        return color.__str__().replace("RGBColor(","").replace(")","").replace(" ","")
    return str(color)

def add_gradient_text_box(slide, text, l, t, w, h,
             size=18, bold=False, color1=C_BLUE, color2=C_BLUE_LIGHT,
             align=PP_ALIGN.LEFT, italic=False):
    """渐变色文字"""
    c1, c2 = _c(color1), _c(color2)
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "PingFang SC"
    # 设置渐变色
    rPr = run._r.get_or_add_rPr()
    # 清除现有颜色
    for old in rPr.findall(qn("a:solidFill")):
        rPr.remove(old)
    gradFill = etree.SubElement(rPr, qn("a:gradFill"), attrib={"rotate": "1"})
    gs = etree.SubElement(gradFill, qn("a:gsLst"))
    for val, pos in [(c1,"0"),(c2,"100000")]:
        g = etree.SubElement(gs, qn("a:gs"), attrib={"pos": pos})
        etree.SubElement(g, qn("a:srgbClr"), attrib={"val": val})
    etree.SubElement(gradFill, qn("a:lin"), attrib={"ang": "5400000", "scaled": "1"})
    return txBox

def add_bottom_bar(slide, color=C_BLUE):
    """底部渐变装饰线"""
    bar = add_rect(slide, 0, 7.38, 13.33, 0.12, fill_color=color)
    set_rect_gradient(bar, color, C_BLUE_LIGHT)

# ── 动画 ──
def add_appear_anim(slide, shape, delay_ms=0):
    """添加淡入动画"""
    sp = shape._element
    # 找到 cNvPr 的 id
    nvCxnSpPr = sp.find(qn("p:nvCxnSpPr"))
    if nvCxnSpPr is None:
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            return
        nvCxnSpPr = etree.SubElement(nvSpPr, qn("p:nvCxnSpPr"))

    # 获取 id
    cNvPr = nvCxnSpPr.getprevious()
    if cNvPr is None:
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
    if cNvPr is None:
        return
    shape_id = cNvPr.get("id", "1")

    # 构建 animation
    anim_xml = f'''<p:anim xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{P}">
      <p:par>
        <p:cTn id="1" dur="500ms" fill="hold">
          <p:stCondLst/>
        </p:cTn>
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="1s" fill="hold">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="{delay_ms}ms"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst/>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                                <p:stCondLst/>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr><p:cTn id="6" dur="1"><p:stCondLst/></p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:animVal>
                                    <p:building Cummulative="1">
                                      <p:waveScale="1.25"/>
                                    </p:building>
                                  </p:animVal>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:childTnLst>
        </p:seq>
      </p:childTnLst>
    </p:par>
    <p:prevCondLst/>
    <p:nextCondLst/>
  </p:par>
</p:anim>'''
    try:
        anim_el = etree.fromstring(anim_xml)
        # 插入到 sp 或 txBody
        sp.append(anim_el)
    except:
        pass

# ══════════════════════════════════════════
# 幻灯片 1: 封面
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, C_BLUE_MID)

# 装饰圆形
c1 = add_rect(s, 10, -1.5, 5, 5, fill_color="0057CC")
set_rect_gradient(c1, "0057CC", "0077FF", 3150000)
c2 = add_rect(s, -2, 5, 5, 4, fill_color="003399")
set_rect_gradient(c2, "003399", "002266", 0)
c3 = add_rect(s, 8, 6, 6, 3, fill_color="0066BB")
set_rect_gradient(c3, "0066BB", "0088DD", 1800000)

# 底部渐变线
add_bottom_bar(s, C_BLUE_LIGHT)

# 装饰亮点
dot = add_rect(s, 1.5, 1.0, 0.15, 0.15, fill_color=C_BLUE_LIGHT)

add_text(s, "第十六届三创赛 · 参赛项目", 0.8, 1.2, 8, 0.4,
         size=10, bold=True, color="90B8E0")

# 主标题
add_gradient_text_box(s, "智水清源", 0.8, 1.7, 9, 1.2,
         size=72, bold=True, color1=C_WHITE, color2="A8D4FF")

# 副标题
add_text(s, "Nereus-S1 智能泳池清洁机器人\n清洁 · 监测 · 投药，一机全搞定",
         0.8, 3.1, 9, 1.0, size=22, color="C0D8F0")

# 三个指标卡片
for i, (num, label) in enumerate([
    ("8L", "垃圾篮容量"),
    ("35%", "作业时间降低"),
    ("3合1", "传感器融合"),
]):
    x = 0.8 + i * 3.0
    card = add_rect(s, x, 4.5, 2.6, 1.2, line_color="0077FF")
    set_rect_gradient(card, "003A88", "0057CC", 0)
    # 顶部强调线
    top_line = add_rect(s, x, 4.5, 2.6, 0.06, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(top_line, C_BLUE_LIGHT, C_BLUE_MID)
    add_text(s, num, x+0.1, 4.6, 2.4, 0.65,
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, x+0.1, 5.2, 2.4, 0.4,
             size=10, color="90B8E0", align=PP_ALIGN.CENTER)

add_text(s, "2026年4月 · 路演", 0.8, 6.7, 5, 0.4,
         size=11, color="6080B0")
add_text(s, "沭波特智能科技（深圳）有限公司",
         7, 6.7, 6, 0.4, size=11, color="6080B0",
         align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════
# 幻灯片 2: 目录
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, "003A80")
add_bottom_bar(s, C_BLUE_LIGHT)

add_text(s, "CONTENTS", 0.8, 0.7, 5, 0.4, size=10, bold=True, color="70A0D0")
add_gradient_text_box(s, "目录", 0.8, 1.1, 8, 0.8,
         size=44, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)

sections = [
    ("01", "痛点", "PAIN POINT"),
    ("02", "方案", "SOLUTION"),
    ("03", "产品", "PRODUCT"),
    ("04", "市场", "MARKET"),
    ("05", "商业模式", "BUSINESS MODEL"),
    ("06", "竞争", "COMPETITION"),
    ("07", "团队", "TEAM"),
    ("08", "规划", "MILESTONE"),
]
cols, rows = 4, 2
for idx, (num, title, en) in enumerate(sections):
    col = idx % cols
    row = idx // cols
    x = 0.8 + col * 3.1
    y = 2.1 + row * 2.5
    card = add_rect(s, x, y, 2.8, 2.0,
             fill_color=None, line_color="1060A0")
    set_rect_gradient(card, "012266", "023A88", 0)
    # 左侧色条
    bar = add_rect(s, x, y, 0.08, 2.0, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(bar, C_BLUE_LIGHT, C_BLUE_MID)
    add_text(s, num, x+0.2, y+0.2, 1.5, 0.5,
             size=28, bold=True, color=C_BLUE_LIGHT)
    add_text(s, title, x+0.2, y+0.8, 2.4, 0.4,
             size=17, bold=True, color=WHITE)
    add_text(s, en, x+0.2, y+1.3, 2.4, 0.4,
             size=9, color="6090C0")

# ══════════════════════════════════════════
# 幻灯片 3: 痛点
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg_blend(s, C_LIGHT_BG, C_WHITE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill_color=C_BLUE_LIGHT)

# 顶部蓝色条
topbar = add_rect(s, 0, 0, 13.33, 0.08, fill_color=C_BLUE)
set_rect_gradient(topbar, C_BLUE_MID, C_BLUE_LIGHT)

add_text(s, "01 · PAIN POINT", 0.6, 0.5, 5, 0.35, size=9, bold=True, color=C_BLUE)
add_gradient_text_box(s, "市场痛点", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)
add_text(s, "传统泳池清洁依赖人工，效率低、成本高、效果差",
         0.6, 1.65, 10, 0.5, size=16, color=TEXT2)

pain_cards = [
    ("💸", "人工成本高", "专业清洁人员月薪 ¥6,000+，频繁清洁费用累积，旺季人手严重不足。", "FF3B30", "FFD0CE"),
    ("📉", "清洁效率低", "现有设备多为有线操作，覆盖不均匀，角落盲区多，维护周期长。", "FF9500", "FFE0B0"),
    ("⏰", "水质监测滞后", "pH/余氯异常靠人工检测，响应慢，夏季高峰易出现水质安全事故。", C_BLUE, C_LIGHT_BG),
]
for i, (emoji, title, desc, color, bg_c) in enumerate(pain_cards):
    x = 0.6 + i * 4.1
    # 卡片
    card = add_rect(s, x, 2.4, 3.9, 4.0, line_color="D0D8E8")
    set_rect_gradient(card, C_WHITE, bg_c)
    # 顶部色块
    top_block = add_rect(s, x, 2.4, 3.9, 0.6, fill_color=color)
    set_rect_gradient(top_block, color, C_BLUE_MID if color==C_BLUE else color)
    add_text(s, emoji, x+0.15, 2.48, 0.5, 0.5, size=24)
    add_text(s, title, x+0.7, 2.5, 3.0, 0.45,
             size=14, bold=True, color=WHITE)
    # 描述
    add_text(s, desc, x+0.15, 3.2, 3.6, 3.0,
             size=12, color=TEXT2)

# ══════════════════════════════════════════
# 幻灯片 4: 解决方案
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, "003A80")
add_bottom_bar(s, C_BLUE_LIGHT)

# 装饰
dec = add_rect(s, 11, 0, 3, 7.5, fill_color="002255")
set_rect_gradient(dec, "002255", "003A88", 1800000)

add_text(s, "02 · SOLUTION", 0.6, 0.5, 5, 0.35, size=9, bold=True, color="70A0D0")
add_gradient_text_box(s, "解决方案", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
add_text(s, "Nereus-S1 一机替代人工，提供全方位智能清洁与实时监测",
         0.6, 1.65, 10, 0.5, size=16, color="C0D8F0")

sol_cards = [
    ("🤖", "全自动化作业", "GPS+惯导+超声波融合导航，一键启动，自动规划全域清洁路径，无需人工干预。"),
    ("💧", "实时水质监测", "pH、余氯、浊度三合一传感器，24小时云端同步，异常实时告警推送手机。"),
    ("💊", "智能精准投药", "装配药剂槽，根据水质数据自动计算投药量，维持pH稳定，减少药剂浪费60%。"),
]
for i, (emoji, title, desc) in enumerate(sol_cards):
    x = 0.6 + i * 4.1
    card = add_rect(s, x, 2.4, 3.9, 3.8, line_color="1060A0")
    set_rect_gradient(card, "012266", "023A88", 0)
    # 顶部强调
    top = add_rect(s, x, 2.4, 3.9, 0.06, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(top, C_BLUE_LIGHT, C_BLUE_MID)
    # 底部强调
    bot = add_rect(s, x, 6.12, 3.9, 0.08, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(bot, C_BLUE_MID, C_BLUE_LIGHT)
    add_text(s, emoji, x+0.15, 2.6, 0.6, 0.6, size=30)
    add_text(s, title, x+0.15, 3.3, 3.6, 0.4,
             size=15, bold=True, color=WHITE)
    add_text(s, desc, x+0.15, 3.85, 3.6, 2.2,
             size=11, color="B0C8E0")

# ══════════════════════════════════════════
# 幻灯片 5: 产品功能进度
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg_blend(s, C_LIGHT_BG, C_WHITE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill_color=C_BLUE_LIGHT)
topbar = add_rect(s, 0, 0, 13.33, 0.08, fill_color=C_BLUE)
set_rect_gradient(topbar, C_BLUE_MID, C_BLUE_LIGHT)

add_text(s, "03 · PRODUCT", 0.6, 0.5, 5, 0.35, size=9, bold=True, color=C_BLUE)
add_gradient_text_box(s, "产品核心功能", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)
add_text(s, "六大核心能力，构建完整泳池管理闭环",
         0.6, 1.65, 10, 0.5, size=16, color=TEXT2)

progs = [
    ("清洁效率", 95),
    ("水质监测覆盖率", 100),
    ("投药精准度", 87),
    ("太阳能自续航", 65),
]
for i, (label, pct) in enumerate(progs):
    y = 2.5 + i * 1.15
    add_text(s, label, 0.6, y, 3, 0.35, size=12, bold=True, color=TEXT)
    add_text(s, f"{pct}%", 10.2, y, 1.5, 0.35,
             size=14, bold=True, color=C_BLUE, align=PP_ALIGN.RIGHT)
    # 进度条底
    bar_bg = add_rect(s, 0.6, y+0.4, 11.1, 0.3, fill_color="E8EEF8")
    # 进度条填充
    fill_w = 11.1 * pct / 100
    bar_fill = add_rect(s, 0.6, y+0.4, fill_w, 0.3, fill_color=C_BLUE)
    set_rect_gradient(bar_fill, C_BLUE_MID, C_BLUE_LIGHT)

# ══════════════════════════════════════════
# 幻灯片 6: 产品规格
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg_blend(s, C_LIGHT_BG, C_WHITE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill_color=C_BLUE_LIGHT)
topbar = add_rect(s, 0, 0, 13.33, 0.08, fill_color=C_BLUE)
set_rect_gradient(topbar, C_BLUE_MID, C_BLUE_LIGHT)

add_text(s, "03 · PRODUCT", 0.6, 0.5, 5, 0.35, size=9, bold=True, color=C_BLUE)
add_gradient_text_box(s, "产品规格", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)

specs = [
    ("8L", "垃圾篮容量"),
    ("13.5W", "太阳能功率"),
    ("35%", "作业时间降低"),
    ("3合1", "传感器融合"),
    ("蓝牙", "APP控制"),
    ("IP68", "防水等级"),
]
for i, (val, label) in enumerate(specs):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.15
    y = 2.0 + row * 2.6
    card = add_rect(s, x, y, 3.9, 2.3, line_color="CCD8FF")
    set_rect_gradient(card, C_WHITE, C_LIGHT_BG)
    # 顶部渐变条
    top = add_rect(s, x, y, 3.9, 0.08, fill_color=C_BLUE)
    set_rect_gradient(top, C_BLUE_MID, C_BLUE_LIGHT)
    add_gradient_text_box(s, val, x+0.15, y+0.3, 3.6, 1.0,
             size=36, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)
    add_text(s, label, x+0.15, y+1.55, 3.6, 0.5,
             size=12, color=TEXT2)

# ══════════════════════════════════════════
# 幻灯片 7: 市场
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, "003A80")
add_bottom_bar(s, C_BLUE_LIGHT)
dec = add_rect(s, 11, 0, 3, 7.5, fill_color="002255")
set_rect_gradient(dec, "002255", "003A88", 1800000)

add_text(s, "04 · MARKET", 0.6, 0.5, 5, 0.35, size=9, bold=True, color="70A0D0")
add_gradient_text_box(s, "市场规模", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
add_text(s, "泳池经济快速增长，智能化升级需求爆发",
         0.6, 1.65, 10, 0.5, size=16, color="C0D8F0")

market_cards = [
    ("¥380亿", "中国泳池市场规模", "2025年", "+18% 年增速", C_GREEN),
    ("12%", "智能清洁设备渗透率", "当前 → 2030年 45%", "+24% CAGR", C_GREEN),
    ("1,200万", "目标用户规模", "B端+C端", "民宿/酒店/私家泳池", C_BLUE_LIGHT),
]
for i, (num, sub, year, trend, trend_color) in enumerate(market_cards):
    x = 0.6 + i * 4.1
    card = add_rect(s, x, 2.3, 3.9, 4.5, line_color="1060A0")
    set_rect_gradient(card, "012266", "023A88", 0)
    top = add_rect(s, x, 2.3, 3.9, 0.08, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(top, C_BLUE_LIGHT, C_BLUE_MID)
    add_text(s, sub, x+0.15, 2.5, 3.6, 0.4,
             size=10, color="6090C0")
    add_gradient_text_box(s, num, x+0.15, 3.0, 3.6, 1.0,
             size=34, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
    add_text(s, year, x+0.15, 4.2, 3.6, 0.35,
             size=10, color="6090C0")
    add_text(s, trend, x+0.15, 4.7, 3.6, 0.5,
             size=13, bold=True, color=rgb(trend_color) if isinstance(trend_color, str) else trend_color)

# ══════════════════════════════════════════
# 幻灯片 8: 商业模式
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg_blend(s, C_LIGHT_BG, C_WHITE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill_color=C_BLUE_LIGHT)
topbar = add_rect(s, 0, 0, 13.33, 0.08, fill_color=C_BLUE)
set_rect_gradient(topbar, C_BLUE_MID, C_BLUE_LIGHT)

add_text(s, "05 · BUSINESS MODEL", 0.6, 0.5, 5, 0.35, size=9, bold=True, color=C_BLUE)
add_gradient_text_box(s, "商业模式", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)
add_text(s, "硬件 + 订阅服务 + 耗材，复合收入结构",
         0.6, 1.65, 10, 0.5, size=16, color=TEXT2)

biz = [
    ("硬件销售", "¥8,800", "Nereus-S1 机器人本体\n一次性销售\n线上+线下渠道", C_BLUE_DARK),
    ("订阅服务", "¥499/年", "云端管家 + AI水质周报\n高复购、高毛利\n现金流稳定", C_BLUE),
    ("上门服务", "¥1,899/季", "深度清洁 + 药剂补给\n面向B端\n高客单价", C_BLUE_MID),
]
for i, (name, price, desc, color) in enumerate(biz):
    x = 0.6 + i * 4.1
    card = add_rect(s, x, 2.3, 3.9, 4.5, line_color="D0D8E8")
    set_rect_gradient(card, C_WHITE, C_LIGHT_BG)
    # 顶部色块
    top = add_rect(s, x, 2.3, 3.9, 0.7, fill_color=color)
    set_rect_gradient(top, color, C_BLUE_MID)
    add_text(s, name, x+0.15, 2.4, 3.6, 0.5,
             size=14, bold=True, color=WHITE)
    add_gradient_text_box(s, price, x+0.15, 3.1, 3.6, 0.8,
             size=30, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)
    add_text(s, desc, x+0.15, 4.0, 3.6, 2.5,
             size=12, color=TEXT2)

# ══════════════════════════════════════════
# 幻灯片 9: 竞争优势
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, "003A80")
add_bottom_bar(s, C_BLUE_LIGHT)

add_text(s, "06 · COMPETITION", 0.6, 0.5, 5, 0.35, size=9, bold=True, color="70A0D0")
add_gradient_text_box(s, "竞争优势", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
add_text(s, "三大壁垒，构建难以复制的护城河",
         0.6, 1.65, 10, 0.5, size=16, color="C0D8F0")

adv = [
    ("💡", "技术壁垒", "多传感器融合算法、边缘计算、自适应路径规划，3项核心专利申请中。"),
    ("📊", "数据壁垒", "水质数据库持续积累，AI模型随用户规模不断优化，后来者难以追赶。"),
    ("🌐", "生态壁垒", "硬件+APP+云平台三位一体，订阅服务锁定用户，高转换成本降低流失率。"),
]
for i, (emoji, title, desc) in enumerate(adv):
    x = 0.6 + i * 4.1
    card = add_rect(s, x, 2.4, 3.9, 4.3, line_color="1060A0")
    set_rect_gradient(card, "012266", "023A88", 0)
    top = add_rect(s, x, 2.4, 3.9, 0.06, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(top, C_BLUE_LIGHT, C_BLUE_MID)
    bot = add_rect(s, x, 6.64, 3.9, 0.06, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(bot, C_BLUE_MID, C_BLUE_LIGHT)
    add_text(s, emoji, x+0.15, 2.6, 0.6, 0.6, size=28)
    add_text(s, title, x+0.15, 3.3, 3.6, 0.5,
             size=16, bold=True, color=WHITE)
    add_text(s, desc, x+0.15, 3.9, 3.6, 2.5,
             size=12, color="B0C8E0")

# ══════════════════════════════════════════
# 幻灯片 10: 团队
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg_blend(s, C_LIGHT_BG, C_WHITE)
add_rect(s, 0, 7.38, 13.33, 0.12, fill_color=C_BLUE_LIGHT)
topbar = add_rect(s, 0, 0, 13.33, 0.08, fill_color=C_BLUE)
set_rect_gradient(topbar, C_BLUE_MID, C_BLUE_LIGHT)

add_text(s, "07 · TEAM", 0.6, 0.5, 5, 0.35, size=9, bold=True, color=C_BLUE)
add_gradient_text_box(s, "核心团队", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=C_BLUE_DARK, color2=C_BLUE)
add_text(s, "跨学科背景，产学研深度融合",
         0.6, 1.65, 10, 0.5, size=16, color=TEXT2)

team = [
    ("张XX", "创始人 & CEO", "XX大学 MBA\n10年智能硬件经验", "张", C_BLUE_DARK),
    ("李XX", "CTO", "XX大学 博士\n机器人算法专家", "李", C_BLUE),
    ("王XX", "COO", "前XX科技运营总监\n精益生产", "王", C_BLUE_MID),
    ("陈XX", "市场总监", "泳池行业8年\n渠道经验", "陈", "0057CC"),
]
for i, (name, role, desc, initial, color) in enumerate(team):
    col = i % 4
    x = 0.6 + col * 3.1
    y = 2.4
    card = add_rect(s, x, y, 2.9, 4.3, line_color="CCD8FF")
    set_rect_gradient(card, C_WHITE, C_LIGHT_BG)
    # 头像渐变圆
    avatar_bg = add_rect(s, x+0.95, y+0.3, 1.0, 1.0, fill_color=color)
    set_rect_gradient(avatar_bg, color, C_BLUE_MID)
    add_text(s, initial, x+0.95, y+0.35, 1.0, 0.9,
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, name, x+0.15, y+1.5, 2.6, 0.5,
             size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, role, x+0.15, y+2.0, 2.6, 0.4,
             size=11, color=C_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.15, y+2.5, 2.6, 2.0,
             size=10, color=TEXT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# 幻灯片 11: 里程碑
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, "003A80")
add_bottom_bar(s, C_BLUE_LIGHT)
dec = add_rect(s, 10.5, 0, 3.5, 7.5, fill_color="002255")
set_rect_gradient(dec, "002255", "003A88", 1800000)

add_text(s, "08 · MILESTONE", 0.6, 0.5, 5, 0.35, size=9, bold=True, color="70A0D0")
add_gradient_text_box(s, "发展规划", 0.6, 0.85, 8, 0.8,
         size=38, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
add_text(s, "清晰路径，稳扎稳打", 0.6, 1.65, 10, 0.5, size=16, color="C0D8F0")

milestones = [
    ("2026 Q2", "完成产品量产", "小批量生产500台，建立供应链"),
    ("2026 Q4", "种子用户验证", "B端民宿30家+C端50台，月流水破10万"),
    ("2027 Q2", "A轮融资", "目标融资500万，用于扩产+市场推广"),
    ("2028", "市场占有率目标", "B端TOP3，年营收突破3000万"),
]
# 时间轴线
axis = add_rect(s, 1.2, 2.4, 0.05, 4.6, fill_color=C_BLUE_LIGHT)
set_rect_gradient(axis, C_BLUE_LIGHT, C_BLUE_MID)
for i, (date, text, desc) in enumerate(milestones):
    y = 2.4 + i * 1.15
    dot = add_rect(s, 1.08, y, 0.3, 0.3, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(dot, C_BLUE_LIGHT, C_BLUE_MID)
    add_text(s, date, 1.55, y-0.02, 1.5, 0.35,
             size=10, bold=True, color=C_BLUE_LIGHT)
    add_text(s, text, 1.55, y+0.32, 4.8, 0.4,
             size=13, bold=True, color=WHITE)
    add_text(s, desc, 1.55, y+0.75, 4.8, 0.4,
             size=10, color="B0C8E0")

# 右侧目标卡
for j, (label, num, sub) in enumerate([
    ("累计融资金额目标", "¥800万", "种子轮 + A轮"),
    ("3年后估值目标", "¥1亿", "B轮前"),
]):
    y = 2.4 + j * 2.4
    card = add_rect(s, 7.5, y, 5.3, 2.1, line_color="1060A0")
    set_rect_gradient(card, "012266", "023A88", 0)
    top = add_rect(s, 7.5, y, 5.3, 0.06, fill_color=C_BLUE_LIGHT)
    set_rect_gradient(top, C_BLUE_LIGHT, C_BLUE_MID)
    add_text(s, label, 7.7, y+0.15, 4.9, 0.4, size=10, color="6090C0")
    add_gradient_text_box(s, num, 7.7, y+0.55, 4.9, 0.9,
             size=34, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
    add_text(s, sub, 7.7, y+1.55, 4.9, 0.35, size=10, color="6090C0")

# ══════════════════════════════════════════
# 幻灯片 12: 感谢页
# ══════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_gradient_bg(s, C_BLUE_DARK, C_BLUE, "003A80")
add_bottom_bar(s, C_BLUE_LIGHT)

# 装饰圆
dc1 = add_rect(s, -1, 4, 6, 5, fill_color="002255")
set_rect_gradient(dc1, "002255", "003A88", 3150000)
dc2 = add_rect(s, 9, -2, 6, 5, fill_color="0057CC")
set_rect_gradient(dc2, "0057CC", "0077FF", 1800000)

add_text(s, "THANKS", 0.6, 0.5, 5, 0.4, size=10, bold=True, color="70A0D0")
add_gradient_text_box(s, "感谢聆听", 0.6, 1.0, 12, 1.0,
         size=60, bold=True, color1=WHITE, color2=C_BLUE_LIGHT)
add_text(s, "期待与您深入交流，共同推动泳池智能化变革",
         0.6, 2.2, 12, 0.5, size=20, color="C0D8F0")

contacts = [
    ("📧  contact@shuboat.com"),
    ("📱  138-0000-0000"),
    ("🌐  www.shuboat.com"),
]
for i, val in enumerate(contacts):
    x = 0.6 + i * 4.1
    card = add_rect(s, x, 3.5, 3.9, 0.6, line_color="1060A0")
    set_rect_gradient(card, "012266", "023A88", 0)
    add_text(s, val, x+0.1, 3.58, 3.7, 0.45, size=14, color=WHITE)

add_text(s, "沭波特智能科技（深圳）有限公司",
         0, 5.5, 13.33, 0.5, size=14, color="6090C0",
         align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════
# 保存
# ══════════════════════════════════════════
out = "/Users/louis/Desktop/三创赛/zhishui-landing/智水清源-路演模版.pptx"
prs.save(out)
print(f"✅ 已保存: {out}")
